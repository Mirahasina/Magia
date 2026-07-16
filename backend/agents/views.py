from .messaging_service import MessagingService
import logging
import os
import re
import datetime
import json
import threading
import time
import traceback
import urllib.parse

import requests
from google import genai
import environ

logger = logging.getLogger(__name__)
import pandas as pd
import PyPDF2
import docx
from accounts.models import User, PLAN_LIMITS, WorkspaceMember
from django.utils import timezone
from rest_framework import viewsets, status, parsers, permissions
from rest_framework.permissions import AllowAny
from rest_framework.decorators import action
from rest_framework.response import Response
from odf.opendocument import load as odf_load
from odf.text import P
from odf import teletype
from pptx import Presentation
from .models import (
    Agent, KnowledgeBase, Template, WhatsAppConfig, ChatMessage, EmailConfig, LinkedInConfig, FacebookConfig, AgentTeam, AgentLink, ContactAssignment, AuditLog, Contact,
    ProspectSearchJob, ProspectLead,
)
from .serializers import (
    AgentSerializer, KnowledgeBaseSerializer, TemplateSerializer, 
    WhatsAppConfigSerializer, ChatMessageSerializer, EmailConfigSerializer,
    AgentFeedbackSerializer, AgentTeamSerializer, AgentLinkSerializer,
    ContactAssignmentSerializer, AuditLogSerializer, LinkedInConfigSerializer,
    FacebookConfigSerializer, ProspectSearchJobSerializer,
)
from .llm_service import get_llm_response, classify_pertinence, DEFAULT_GEMINI_MODELS
from .email_service import (
    process_emails_for_agent, sync_email_history,
    send_email_reply, refresh_google_token,
)
from django.db.models import Q, Count, Max as MaxAgg
from django.db.models.functions import TruncDay
from django.shortcuts import redirect
from rest_framework.exceptions import PermissionDenied
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from magia_backend.settings_constants import FRONTEND_URL
from .rag_service import add_texts_to_knowledge_base, search_knowledge_base, search_agent_and_team_knowledge_base
from .prospection_service import (
    mark_prospect_replied,
    purge_passive_crm_contacts,
    schedule_followup_after_ai,
    infer_channels_for_team_agent,
    attach_user_channel_configs,
    build_handoff_intro,
)



env = environ.Env()
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
environ.Env.read_env(os.path.join(BASE_DIR, '.env'))



def sync_facebook_history(config, user):
    try:
        resp = requests.get(
            f"https://graph.facebook.com/v18.0/{config.page_id}/conversations",
            params={
                'access_token': config.page_access_token,
                'fields': 'participants{id,name},messages{message,from,created_time,id}',
            },
            timeout=15,
        )
        if resp.status_code != 200:
            err = resp.json().get('error', {}).get('message', 'Erreur Graph API')
            return {'imported': 0, 'error': err, 'status_code': 400}
        data = resp.json()
        synced = 0
        for conv in data.get('data', []):
            participants = (conv.get('participants') or {}).get('data') or []
            peer_id = None
            peer_name = None
            for p in participants:
                pid = p.get('id')
                if pid and str(pid) != str(config.page_id):
                    peer_id = str(pid)
                    peer_name = p.get('name') or peer_id
                    break
            if not peer_id:
                continue
            for msg_data in (conv.get('messages') or {}).get('data') or []:
                mid = msg_data.get('id', '')
                text = (msg_data.get('message') or '').strip()
                sender_id = str((msg_data.get('from') or {}).get('id', '') or '')
                is_me = sender_id == str(config.page_id)
                if not text:
                    continue
                if ChatMessage.objects.filter(user=user, whatsapp_message_id=mid).exists():
                    continue
                ChatMessage.objects.create(
                    user=user,
                    sender='ai' if is_me else 'user',
                    content=text,
                    contact_info=peer_id,
                    contact_name=peer_name,
                    source='facebook',
                    whatsapp_message_id=mid,
                    is_read=True,
                    status='archived',
                )
                synced += 1
        return {'imported': synced, 'error': None}
    except Exception as exc:
        return {'imported': 0, 'error': str(exc), 'status_code': 500}


def facebook_auto_reply(user_id, agent_id, sender_psid, text):
    try:
        user = User.objects.get(id=user_id)
        agent = Agent.objects.get(id=agent_id)
        if not agent.is_active:
            return

        history = list(
            ChatMessage.objects.filter(
                user=user, contact_info=sender_psid, source='facebook'
            ).order_by('-created_at')[:6]
        )
        history.reverse()
        history_str = "\n".join(f"{m.sender}: {m.content}" for m in history)

        rag_context = search_knowledge_base(agent.id, text, top_k=4)
        context_for_llm = (
            "Documents sources (Extraits pertinents par recherche sémantique RAG) :\n"
            f"{rag_context if rag_context.strip() else '[Aucun document fourni]'}\n\n"
            f"Historique récent :\n{history_str}"
        )

        answer = get_llm_response(
            agent_name=agent.name,
            agent_role=agent.role,
            system_prompt=agent.system_prompt,
            knowledge_context=context_for_llm,
            user_message=text,
            model_name=agent.llm_model,
        )

        sent = MessagingService.send_message(user, sender_psid, answer, 'facebook')
        if not sent:
            logger.warning("Facebook auto-reply: envoi échoué pour PSID %s", sender_psid)
            return

        ChatMessage.objects.create(
            user=user,
            agent=agent,
            sender='ai',
            content=answer,
            contact_info=sender_psid,
            source='facebook',
            status='new',
        )
        schedule_followup_after_ai(user, sender_psid, 'facebook', analyze=False)
    except Exception as exc:
        logger.error("Facebook auto-reply error (PSID %s): %s", sender_psid, exc)


def get_whatsapp_port(user_id):
    if user_id == 'global':
        return 3001
    uid_str = str(user_id)
    return 3001 + (sum(ord(c) for c in uid_str) % 100)

def check_ai_quota(user):
    plan = 'gratuit'
    if hasattr(user, 'subscription'):
        plan = user.subscription.plan_name
    
    if plan == 'entreprise':
        return True, None
        
    limit = 50 if plan == 'gratuit' else 1000
    last_24h = timezone.now() - datetime.timedelta(hours=24)
    count = ChatMessage.objects.filter(user=user, sender='ai', created_at__gt=last_24h).count()
    
    if count >= limit:
        return False, f"Quota journalier atteint pour le plan {plan} ({limit} messages). Veuillez passer au plan supérieur."
    return True, None


def _extract_text(file_obj) -> str:
    filename = file_obj.name.lower()
    if filename.endswith('.txt'):
        return file_obj.read().decode('utf-8')
    if filename.endswith('.pdf'):
        reader = PyPDF2.PdfReader(file_obj)
        return ''.join(
            (page.extract_text() or '') + '\n'
            for page in reader.pages
        )
    if filename.endswith('.docx'):
        doc = docx.Document(file_obj)
        return '\n'.join(para.text for para in doc.paragraphs)
    if filename.endswith(('.xlsx', '.xls')):
        return pd.read_excel(file_obj).to_string()
    if filename.endswith('.pptx'):
        prs = Presentation(file_obj)
        return '\n'.join(
            getattr(shape, 'text', '')
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, 'text', None) is not None
        )
    if filename.endswith(('.odt', '.ods', '.odp')):
        doc = odf_load(file_obj)
        paragraphs = doc.text.getElementsByType(P)
        return '\n'.join(teletype.extractText(p) for p in paragraphs)
    return ''

class WhatsAppConfigViewSet(viewsets.ModelViewSet):
    serializer_class = WhatsAppConfigSerializer

    def get_queryset(self):
        return WhatsAppConfig.objects.filter(user=self.request.user)

    def perform_create(self, serializer):
        serializer.save(user=self.request.user)

    def _auth_token_for_wa(self, request):
        auth = request.META.get('HTTP_AUTHORIZATION', '')
        if auth.lower().startswith('bearer '):
            return auth.split(' ', 1)[1].strip()
        user = request.user
        if getattr(user, 'master_api_key', None):
            return user.master_api_key
        return ''

    @action(detail=True, methods=['post'])
    def start_session(self, request, pk=None):
        """
        Start Baileys and wait for a QR code.
        If the account is not currently connected, credentials are wiped first so
        a fresh QR scan is always required (no silent auto-reconnect).
        """
        from .whatsapp_process import (
            start_for_user, wait_for_qr, is_running, clear_auth_for_user, stop_for_user,
        )

        config = self.get_object()

        if config.is_connected:
            # Already linked - do not wipe session
            if not is_running(request.user.id):
                start_for_user(
                    request.user,
                    auth_token=self._auth_token_for_wa(request),
                    force_restart=False,
                )
            return Response({
                'status': 'connected',
                'phone_number': config.phone_number,
                'service_running': is_running(request.user.id),
            })

        # Security: disconnected -> must scan again (no silent session restore)
        stop_for_user(request.user.id)
        clear_auth_for_user(request.user.id)
        config.is_connected = False
        config.qr_code = None
        config.phone_number = None
        config.save(update_fields=['is_connected', 'qr_code', 'phone_number'])

        result = start_for_user(
            request.user,
            auth_token=self._auth_token_for_wa(request),
            force_restart=True,
        )
        if result.get('error') and not result.get('already_running'):
            return Response(
                {'error': result['error'], 'status': 'error'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR,
            )

        qr = wait_for_qr(request.user, timeout_seconds=20.0)
        config.refresh_from_db()

        # Refuse "connected without QR" after a wipe - must have scanned
        if config.is_connected and not (qr or config.qr_code):
            # Race: session somehow restored - wipe again and ask for QR
            stop_for_user(request.user.id)
            clear_auth_for_user(request.user.id)
            config.is_connected = False
            config.qr_code = None
            config.save(update_fields=['is_connected', 'qr_code'])
            start_for_user(
                request.user,
                auth_token=self._auth_token_for_wa(request),
                force_restart=True,
            )
            qr = wait_for_qr(request.user, timeout_seconds=15.0)
            config.refresh_from_db()

        if qr or config.qr_code:
            return Response({
                'status': 'qr_ready',
                'qr_code': qr or config.qr_code,
                'service_running': True,
                'started': result.get('started'),
                'already_running': result.get('already_running'),
            })

        return Response({
            'status': 'generating_qr',
            'message': (
                'Service WhatsApp démarré. Scannez le QR code qui apparaît '
                '(aucune connexion automatique sans scan).'
            ),
            'service_running': is_running(request.user.id),
            'started': result.get('started'),
            'already_running': result.get('already_running'),
        })

    @action(detail=True, methods=['post'])
    def stop_session(self, request, pk=None):
        from .whatsapp_process import disconnect_for_user
        config = self.get_object()
        disconnect_for_user(request.user.id)
        config.is_connected = False
        config.qr_code = None
        config.phone_number = None
        config.save(update_fields=['is_connected', 'qr_code', 'phone_number'])
        return Response({'status': 'disconnected'})

    @action(detail=True, methods=['get'])
    def get_connection_url(self, request, pk=None):
        """
        Poll status / current QR. Does NOT wipe credentials (use start_session for that).
        """
        from .whatsapp_process import start_for_user, is_running, wait_for_qr

        config = self.get_object()
        if config.is_connected:
            return Response({
                "status": "connected",
                "phone_number": config.phone_number,
            })

        if config.qr_code:
            return Response({
                "qr_code": config.qr_code,
                "status": "qr_ready",
            })

        # Keep process alive while user waits for QR (after start_session)
        if not is_running(request.user.id):
            start_result = start_for_user(
                request.user,
                auth_token=self._auth_token_for_wa(request),
                force_restart=False,
            )
            if start_result.get('error') and not start_result.get('already_running'):
                return Response({
                    "status": "error",
                    "error": start_result['error'],
                    "message": start_result['error'],
                }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        qr = wait_for_qr(request.user, timeout_seconds=8.0)
        config.refresh_from_db()
        if config.is_connected:
            return Response({
                "status": "connected",
                "phone_number": config.phone_number,
            })
        if qr or config.qr_code:
            return Response({
                "qr_code": qr or config.qr_code,
                "status": "qr_ready",
            })
        return Response({
            "status": "generating_qr",
            "message": (
                "Le service WhatsApp génère votre QR code. "
                "Scannez-le pour connecter - aucune connexion automatique."
            ),
        })

    @action(detail=True, methods=['get'])
    def refresh_connection(self, request, pk=None):
        from .whatsapp_process import auth_dir_for_user

        config = self.get_object()
        creds = auth_dir_for_user(request.user.id) / 'creds.json'
        # Stale DB flag without session file => treat as disconnected (must rescan)
        if config.is_connected and not creds.exists():
            config.is_connected = False
            config.phone_number = None
            config.qr_code = None
            config.save(update_fields=['is_connected', 'phone_number', 'qr_code'])
        return Response({
            "status": "connected" if config.is_connected else "not_connected",
            "is_connected": config.is_connected,
            "phone_number": config.phone_number,
            "qr_code": config.qr_code
        })

    @action(detail=True, methods=['post'])
    def sync_messages(self, request, pk=None):
        return Response({"status": "Synchronisation en temps réel via le service local active."})

    @action(detail=False, methods=['post'], permission_classes=[AllowAny], authentication_classes=[])
    def update_qr(self, request):
        user_id = request.data.get('user_id')
        qr_code = request.data.get('qr_code')
        if not user_id or user_id == 'global':
            return Response({'error': 'User ID required'}, status=400)
        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)
        
        config = WhatsAppConfig.objects.filter(user=user).order_by('-updated_at').first()
        if not config:
            return Response({'status': 'ignored', 'message': 'No WhatsApp config for user'})
        config.qr_code = qr_code
        config.is_connected = False
        config.save(update_fields=['qr_code', 'is_connected', 'updated_at'])
        return Response({'status': 'updated'})

    @action(detail=False, methods=['post'], permission_classes=[AllowAny], authentication_classes=[])
    def set_connected(self, request):
        user_id = request.data.get('user_id')
        phone_number = request.data.get('phone_number')
        if not user_id or user_id == 'global':
            return Response({'error': 'User ID required'}, status=400)
        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)

        # Never auto-create a default account - user must add one in Paramètres first
        config = WhatsAppConfig.objects.filter(user=user).order_by('-updated_at').first()
        if not config:
            return Response({'status': 'ignored', 'message': 'No WhatsApp config for user'})
        config.is_connected = True
        config.phone_number = phone_number
        config.qr_code = None
        config.save(update_fields=['is_connected', 'phone_number', 'qr_code', 'updated_at'])
        return Response({'status': 'connected'})

    @action(detail=False, methods=['post'], permission_classes=[AllowAny], authentication_classes=[])
    def set_disconnected(self, request):
        user_id = request.data.get('user_id')
        if not user_id or user_id == 'global':
            return Response({'error': 'User ID required'}, status=400)
        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)

        from .whatsapp_process import disconnect_for_user
        disconnect_for_user(user.id)

        configs = WhatsAppConfig.objects.filter(user=user)
        configs.update(is_connected=False, qr_code=None, phone_number=None)
        return Response({'status': 'disconnected'})

    @action(detail=False, methods=['post'], permission_classes=[AllowAny], authentication_classes=[])
    def sync_whatsapp_contacts(self, request):
        user_id = request.data.get('user_id')
        contacts = request.data.get('contacts', [])
        if not user_id or user_id == 'global':
            return Response({'error': 'User ID required'}, status=400)
        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)
        
        # Address-book sync must NOT flood the CRM pipeline. Only refresh names
        # for prospects already intentionally added to Prospection.
        updated = 0
        for c in contacts:
            contact_info = c.get('id')
            if not contact_info:
                continue
            name = c.get('name') or c.get('notify') or c.get('verifiedName') or contact_info.split('@')[0]
            contact = Contact.objects.filter(
                user=user, source='whatsapp', contact_info=contact_info
            ).first()
            if contact and name and contact.name != name:
                contact.name = name
                contact.save(update_fields=['name'])
                updated += 1

        purged = purge_passive_crm_contacts(user)
        return Response({
            'status': 'success',
            'synced': 0,
            'updated_existing': updated,
            'purged_passive': purged,
        })

    def perform_destroy(self, instance):
        from .whatsapp_process import disconnect_for_user
        if instance.user_id:
            disconnect_for_user(instance.user_id)
        ChatMessage.objects.filter(user=instance.user, source='whatsapp').delete()
        Contact.objects.filter(user=instance.user, source='whatsapp').delete()
        super().perform_destroy(instance)

    @action(detail=False, methods=['post'], permission_classes=[AllowAny], authentication_classes=[])
    def whatsapp_gateway(self, request):
        user_id = request.data.get('user_id')
        message = request.data.get('message', '').strip()
        wa_id = request.data.get('message_id', '')
        wa_sender = request.data.get('sender', '')
        push_name = request.data.get('push_name')
        is_historical = request.data.get('is_historical', False)
        is_me = request.data.get('is_me', False)
        
        if not user_id or user_id == 'global':
            return Response({'error': 'User ID required'}, status=400)
            
        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)

        if wa_id and ChatMessage.objects.filter(user=user, whatsapp_message_id=wa_id).exists():
            return Response({'status': 'ignored', 'reason': 'duplicate'})

        assignment = ContactAssignment.objects.filter(user=user, contact_info=wa_sender).first()
        effective_agent = None
        if assignment:
            effective_agent = assignment.agent
        else:
            active_agents = Agent.objects.filter(user=user, is_active=True)
            for agent in active_agents:
                if agent.channels and any(str(c).lower() == 'whatsapp' for c in agent.channels):
                    effective_agent = agent
                    break
            if not effective_agent:
                effective_agent = active_agents.first()

        ChatMessage.objects.create(
            user=user,
            agent=effective_agent,
            sender='ai' if is_me else 'user',
            content=message,
            contact_info=wa_sender,
            contact_name=push_name,
            source='whatsapp',
            whatsapp_message_id=wa_id,
            is_read=bool(is_historical),
            status='archived' if is_historical else 'new'
        )

        if not is_historical and message:
            if is_me:
                schedule_followup_after_ai(user, wa_sender, 'whatsapp', analyze=False)
            else:
                mark_prospect_replied(user, wa_sender, 'whatsapp')

        return Response({'status': 'received'})


    @action(detail=False, methods=['post'], permission_classes=[AllowAny], authentication_classes=[])
    def whatsapp_gateway_bulk(self, request):
        user_id = request.data.get('user_id')
        messages = request.data.get('messages', [])
        
        if not user_id or user_id == 'global':
            return Response({'error': 'User ID required'}, status=400)
            
        try:
            user = User.objects.get(id=user_id)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=404)

        existing_ids = set(ChatMessage.objects.filter(user=user, source='whatsapp').values_list('whatsapp_message_id', flat=True))
        
        new_messages = []
        contacts_to_update = {}
        
        for msg_data in messages:
            wa_id = msg_data.get('message_id', '')
            if not wa_id or wa_id in existing_ids:
                continue
                
            wa_sender = msg_data.get('sender', '')
            push_name = msg_data.get('push_name')
            is_me = msg_data.get('is_me', False)
            content = msg_data.get('message', '').strip()
            
            if not content: continue
            
            new_messages.append(ChatMessage(
                user=user,
                sender='ai' if is_me else 'user',
                content=content,
                contact_info=wa_sender,
                contact_name=push_name,
                source='whatsapp',
                whatsapp_message_id=wa_id,
                is_read=True,
                status='archived'
            ))
            
            if wa_sender not in contacts_to_update or push_name:
                contacts_to_update[wa_sender] = push_name
        
        if new_messages:
            ChatMessage.objects.bulk_create(new_messages, ignore_conflicts=True)
            # Update display names on existing CRM prospects only - never auto-create.
            for sender, name in contacts_to_update.items():
                if not name:
                    continue
                Contact.objects.filter(
                    user=user, source='whatsapp', contact_info=sender
                ).exclude(name=name).update(name=name)

        return Response({'status': 'success', 'processed': len(new_messages)})


class EmailConfigViewSet(viewsets.ModelViewSet):
    serializer_class = EmailConfigSerializer

    def get_queryset(self):
        # oauth2_callback is AllowAny and has no authenticated user
        if getattr(self, 'action', None) == 'oauth2_callback':
            return EmailConfig.objects.all()
        return EmailConfig.objects.filter(user=self.request.user)

    def perform_create(self, serializer):
        serializer.save(user=self.request.user)

    def perform_update(self, serializer):
        """Allow updating imap_server, smtp_server, email, password."""
        instance = serializer.save()
        # Save password separately (not in serializer for security, sent via PATCH)
        password = self.request.data.get('password')
        if password:
            instance.password = password
            instance.save(update_fields=['password'])

    def _google_oauth_credentials(self):
        # Strip whitespace/quotes - common cause of Google "invalid_client"
        client_id = env('GOOGLE_CLIENT_ID', default='').strip().strip('"').strip("'")
        client_secret = env('GOOGLE_CLIENT_SECRET', default='').strip().strip('"').strip("'")
        redirect_uri = env('GOOGLE_REDIRECT_URI', default='').strip().strip('"').strip("'")
        return client_id, client_secret, redirect_uri

    def _frontend_gmail_redirect(self, **params):
        query = urllib.parse.urlencode(params)
        base = FRONTEND_URL.rstrip('/')
        return f"{base}/?view=integration&{query}"

    @action(detail=True, methods=['get'])
    def get_connection_url(self, request, pk=None):
        """Build Google OAuth URL for one-click Gmail connection."""
        self.get_object()  # ensure config belongs to the current user
        client_id, client_secret, redirect_uri = self._google_oauth_credentials()
        if not client_id or not client_secret or not redirect_uri:
            return Response(
                {
                    "error": (
                        "GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET ou GOOGLE_REDIRECT_URI "
                        "n'est pas configuré dans le fichier .env du backend."
                    )
                },
                status=status.HTTP_400_BAD_REQUEST,
            )

        params = {
            'client_id': client_id,
            'redirect_uri': redirect_uri,
            'response_type': 'code',
            'scope': ' '.join([
                'openid',
                'email',
                'profile',
                'https://mail.google.com/',
            ]),
            'access_type': 'offline',
            'prompt': 'consent',
            'state': str(pk),
        }
        url = f"https://accounts.google.com/o/oauth2/v2/auth?{urllib.parse.urlencode(params)}"
        return Response({"url": url})

    @action(
        detail=False,
        methods=['get'],
        url_path='oauth2_callback',
        permission_classes=[AllowAny],
        authentication_classes=[],
    )
    def oauth2_callback(self, request):
        """Google redirects here after consent; store tokens and send user back to the app."""
        error = request.query_params.get('error')
        if error:
            return redirect(self._frontend_gmail_redirect(gmail_error=error))

        code = request.query_params.get('code')
        state = request.query_params.get('state')
        if not code or not state:
            return redirect(self._frontend_gmail_redirect(gmail_error='missing_code'))

        try:
            config_id = int(state)
        except (TypeError, ValueError):
            return redirect(self._frontend_gmail_redirect(gmail_error='invalid_state'))

        try:
            config = EmailConfig.objects.get(pk=config_id)
        except EmailConfig.DoesNotExist:
            return redirect(self._frontend_gmail_redirect(gmail_error='config_not_found'))

        client_id, client_secret, redirect_uri = self._google_oauth_credentials()
        if not client_id or not client_secret or not redirect_uri:
            return redirect(self._frontend_gmail_redirect(gmail_error='google_not_configured'))

        try:
            token_resp = requests.post(
                'https://oauth2.googleapis.com/token',
                data={
                    'code': code,
                    'client_id': client_id,
                    'client_secret': client_secret,
                    'redirect_uri': redirect_uri,
                    'grant_type': 'authorization_code',
                },
                timeout=15,
            )
            if token_resp.status_code != 200:
                logger.error('Google token exchange failed: %s', token_resp.text)
                return redirect(self._frontend_gmail_redirect(gmail_error='token_exchange_failed'))

            tokens = token_resp.json()
            access_token = tokens.get('access_token')
            refresh_token = tokens.get('refresh_token')
            if not access_token:
                return redirect(self._frontend_gmail_redirect(gmail_error='no_access_token'))

            userinfo_resp = requests.get(
                'https://www.googleapis.com/oauth2/v3/userinfo',
                headers={'Authorization': f'Bearer {access_token}'},
                timeout=10,
            )
            if userinfo_resp.status_code != 200:
                return redirect(self._frontend_gmail_redirect(gmail_error='userinfo_failed'))

            email = userinfo_resp.json().get('email')
            if not email:
                return redirect(self._frontend_gmail_redirect(gmail_error='email_missing'))

            config.email = email
            config.oauth_token = access_token
            if refresh_token:
                config.refresh_token = refresh_token
            config.imap_server = 'imap.gmail.com'
            config.smtp_server = 'smtp.gmail.com'
            config.password = ''
            config.is_active = True
            if not config.name or config.name in ('Default Email', 'Nouveau Email'):
                config.name = email
            config.save()

            thread = threading.Thread(target=sync_email_history, args=(config,))
            thread.daemon = True
            thread.start()

            return redirect(self._frontend_gmail_redirect(gmail_connected='1'))
        except Exception as exc:
            logger.exception('Gmail OAuth callback failed: %s', exc)
            return redirect(self._frontend_gmail_redirect(gmail_error='callback_failed'))

    @action(detail=True, methods=['get'])
    def refresh_connection(self, request, pk=None):
        """Verify / refresh Gmail OAuth token. Google OAuth only (no SMTP password)."""
        config = self.get_object()
        if not config.oauth_token and not config.refresh_token:
            if config.is_active:
                config.is_active = False
                config.save(update_fields=['is_active'])
            return Response({
                'status': 'disconnected',
                'is_active': False,
                'email': config.email,
                'message': 'Seule la connexion Google est autorisée. Connectez Gmail via OAuth.',
            })

        token = config.oauth_token
        if config.refresh_token:
            refreshed = refresh_google_token(config)
            if refreshed:
                token = refreshed
                config.refresh_from_db()
        if token:
            try:
                resp = requests.get(
                    'https://www.googleapis.com/oauth2/v3/userinfo',
                    headers={'Authorization': f'Bearer {token}'},
                    timeout=5,
                )
                if resp.status_code == 200:
                    if not config.is_active:
                        config.is_active = True
                        config.save(update_fields=['is_active'])
                    return Response({
                        'status': 'connected',
                        'is_active': True,
                        'email': config.email,
                    })
            except Exception:
                pass
        config.is_active = False
        config.save(update_fields=['is_active'])
        return Response({
            'status': 'disconnected',
            'is_active': False,
            'email': config.email,
            'message': 'Token Google invalide ou expiré. Reconnectez Gmail.',
        })

    @action(detail=True, methods=['post'])
    def test_connection(self, request, pk=None):
        """Google OAuth only - SMTP password login is disabled."""
        return Response(
            {
                'success': False,
                'error': (
                    'La connexion email par SMTP / mot de passe n\'est plus disponible. '
                    'Utilisez « Se connecter avec Google ».'
                ),
            },
            status=status.HTTP_400_BAD_REQUEST,
        )

    @action(detail=True, methods=['post'])
    def configure(self, request, pk=None):
        """SMTP/password activation disabled - use Google OAuth only."""
        return Response(
            {
                'error': (
                    'La configuration SMTP n\'est plus autorisée. '
                    'Connectez Gmail uniquement via Google OAuth.'
                ),
            },
            status=status.HTTP_400_BAD_REQUEST,
        )

    @action(detail=True, methods=['post'])
    def sync_messages(self, request, pk=None):
        """Re-import IMAP history into the inbox. wait=true (default) runs sync in-request."""
        config = self.get_object()
        wait = str(request.data.get('wait', 'true')).lower() not in ('0', 'false', 'no')
        if wait:
            stats = sync_email_history(config)
            if stats.get('error') and stats.get('imported', 0) == 0:
                return Response(
                    {'status': 'error', 'error': stats['error'], **stats},
                    status=status.HTTP_400_BAD_REQUEST,
                )
            return Response({'status': 'done', **stats})
        thread = threading.Thread(target=sync_email_history, args=(config,))
        thread.daemon = True
        thread.start()
        return Response({"status": "Synchronisation Email démarrée"})

    def perform_destroy(self, instance):
        ChatMessage.objects.filter(email_config=instance).delete()
        super().perform_destroy(instance)

@method_decorator(csrf_exempt, name='dispatch')
class FacebookConfigViewSet(viewsets.ModelViewSet):
    serializer_class = FacebookConfigSerializer

    def get_queryset(self):
        if getattr(self, 'action', None) == 'webhook':
            return FacebookConfig.objects.all()
        return FacebookConfig.objects.filter(user=self.request.user)

    def get_permissions(self):
        if getattr(self, 'action', None) == 'webhook':
            return [AllowAny()]
        return super().get_permissions()

    def get_authenticators(self):
        if getattr(self, 'action', None) == 'webhook':
            return []
        return super().get_authenticators()

    def perform_create(self, serializer):
        serializer.save(user=self.request.user)

    @action(detail=True, methods=['get'])
    def get_connection_url(self, request, pk=None):
        fb_app_id = env('FACEBOOK_APP_ID', default='')
        if not fb_app_id:
            return Response({"error": "FACEBOOK_APP_ID n'est pas configuré dans le fichier .env du backend."}, status=status.HTTP_400_BAD_REQUEST)
        
        redirect_uri = request.query_params.get('redirect_uri')
        if not redirect_uri:
            return Response({"error": "redirect_uri est requis."}, status=status.HTTP_400_BAD_REQUEST)
        
        params = {
            'client_id': fb_app_id,
            'redirect_uri': redirect_uri,
            'scope': 'pages_messaging,pages_show_list,pages_read_engagement',
            'state': str(pk)
        }
        url = f"https://www.facebook.com/v18.0/dialog/oauth?{urllib.parse.urlencode(params)}"
        return Response({"url": url})

    @action(detail=True, methods=['post'])
    def exchange_code(self, request, pk=None):
        """Exchange auth code for long-lived user access token and list available Pages."""
        code = request.data.get('code')
        redirect_uri = request.data.get('redirect_uri')
        if not code or not redirect_uri:
            return Response({"error": "code et redirect_uri sont requis."}, status=status.HTTP_400_BAD_REQUEST)
        
        fb_app_id = env('FACEBOOK_APP_ID', default='')
        fb_app_secret = env('FACEBOOK_APP_SECRET', default='')
        if not fb_app_id or not fb_app_secret:
            return Response({"error": "FACEBOOK_APP_ID ou FACEBOOK_APP_SECRET non configuré."}, status=status.HTTP_400_BAD_REQUEST)
        
        # 1. Exchange code for short-lived user access token
        try:
            resp = requests.get(
                "https://graph.facebook.com/v18.0/oauth/access_token",
                params={
                    'client_id': fb_app_id,
                    'redirect_uri': redirect_uri,
                    'client_secret': fb_app_secret,
                    'code': code
                },
                timeout=10
            )
            if resp.status_code != 200:
                err = resp.json().get('error', {}).get('message', 'Impossible d\'échanger le code.')
                return Response({"error": f"Erreur Facebook : {err}"}, status=status.HTTP_400_BAD_REQUEST)
            
            user_token = resp.json().get('access_token')
        except Exception as e:
            return Response({"error": f"Impossible de contacter Facebook : {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        
        # 2. Exchange short-lived user access token for a long-lived user access token
        try:
            resp = requests.get(
                "https://graph.facebook.com/v18.0/oauth/access_token",
                params={
                    'grant_type': 'fb_exchange_token',
                    'client_id': fb_app_id,
                    'client_secret': fb_app_secret,
                    'fb_exchange_token': user_token
                },
                timeout=10
            )
            if resp.status_code == 200:
                user_token = resp.json().get('access_token', user_token)
        except Exception:
            pass
        
        # 3. Retrieve list of pages managed by this user
        try:
            resp = requests.get(
                "https://graph.facebook.com/v18.0/me/accounts",
                params={
                    'access_token': user_token,
                    'fields': 'id,name,access_token'
                },
                timeout=10
            )
            if resp.status_code != 200:
                err = resp.json().get('error', {}).get('message', 'Impossible de récupérer la liste des pages.')
                return Response({"error": f"Erreur Facebook Pages : {err}"}, status=status.HTTP_400_BAD_REQUEST)
            
            pages_data = resp.json().get('data', [])
            pages = []
            for p in pages_data:
                pages.append({
                    'id': p.get('id'),
                    'name': p.get('name'),
                    'access_token': p.get('access_token')
                })
            
            return Response({"pages": pages})
        except Exception as e:
            return Response({"error": f"Erreur de communication Pages : {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    @action(detail=True, methods=['get'])
    def refresh_connection(self, request, pk=None):
        config = self.get_object()
        if config.is_connected and config.page_access_token:
            # Verify token is still valid via Graph API
            try:
                resp = requests.get(
                    f"https://graph.facebook.com/v18.0/{config.page_id}",
                    params={'access_token': config.page_access_token, 'fields': 'id,name'},
                    timeout=5
                )
                if resp.status_code == 200:
                    data = resp.json()
                    config.page_name = data.get('name', config.page_name)
                    config.save(update_fields=['page_name'])
                    return Response({'status': 'connected', 'page_name': config.page_name})
            except Exception:
                pass
            config.is_connected = False
            config.save(update_fields=['is_connected'])
        return Response({'status': 'disconnected', 'message': 'Token invalide ou expiré.'})

    @action(detail=True, methods=['post'])
    def configure(self, request, pk=None):
        """Save Page ID + Page Access Token and verify via Graph API."""
        config = self.get_object()
        page_id = request.data.get('page_id', '').strip()
        page_access_token = request.data.get('page_access_token', '').strip()

        if not page_id or not page_access_token:
            return Response({'error': 'page_id et page_access_token sont requis.'}, status=400)

        # Verify token against Graph API
        try:
            resp = requests.get(
                f"https://graph.facebook.com/v18.0/{page_id}",
                params={'access_token': page_access_token, 'fields': 'id,name'},
                timeout=10
            )
            if resp.status_code != 200:
                err = resp.json().get('error', {}).get('message', 'Token invalide.')
                return Response({'error': f'Échec vérification Graph API : {err}'}, status=400)
            page_data = resp.json()
        except Exception as e:
            return Response({'error': f'Impossible de joindre Facebook : {str(e)}'}, status=500)

        config.page_id = page_id
        config.page_access_token = page_access_token
        config.page_name = page_data.get('name', '')
        config.is_connected = True
        config.save()

        # Subscribe the Page to the app webhook so Messenger events arrive
        subscribe_warning = None
        try:
            sub = requests.post(
                f"https://graph.facebook.com/v18.0/{page_id}/subscribed_apps",
                params={'access_token': page_access_token},
                data={
                    'subscribed_fields': (
                        'messages,messaging_postbacks,messaging_optins,'
                        'message_deliveries,message_reads'
                    ),
                },
                timeout=12,
            )
            if sub.status_code != 200 or not sub.json().get('success'):
                subscribe_warning = (
                    sub.json().get('error', {}).get('message')
                    or f'Abonnement webhook échoué (HTTP {sub.status_code})'
                )
                logger.warning("Facebook subscribed_apps failed for page %s: %s", page_id, subscribe_warning)
        except Exception as exc:
            subscribe_warning = str(exc)
            logger.warning("Facebook subscribed_apps error for page %s: %s", page_id, exc)

        payload = {
            'status': 'connected',
            'page_name': config.page_name,
            'page_id': config.page_id,
        }
        if subscribe_warning:
            payload['webhook_warning'] = (
                f"Page connectée, mais webhook non abonné : {subscribe_warning}. "
                "Les messages live peuvent ne pas arriver tant que le webhook Meta n'est pas configuré."
            )
        return Response(payload)

    @action(detail=True, methods=['post'])
    def send_facebook_message(self, request, pk=None):
        """Send a message to a Facebook user via Graph API (Page Messaging)."""
        config = self.get_object()
        if not config.is_connected:
            return Response({'error': 'Page Facebook non connectée.'}, status=400)
        recipient_id = request.data.get('recipient_psid')
        text = request.data.get('text', '')
        if not recipient_id or not text:
            return Response({'error': 'recipient_psid et text requis.'}, status=400)
        try:
            resp = requests.post(
                f"https://graph.facebook.com/v18.0/{config.page_id}/messages",
                params={'access_token': config.page_access_token},
                json={'recipient': {'id': recipient_id}, 'message': {'text': text}},
                timeout=10
            )
            if resp.status_code == 200:
                return Response({'status': 'sent'})
            return Response({'error': resp.json().get('error', {}).get('message', 'Erreur envoi.')}, status=400)
        except Exception as e:
            return Response({'error': str(e)}, status=500)

    @action(detail=False, methods=['get', 'post'], permission_classes=[AllowAny], authentication_classes=[])
    def webhook(self, request):
        """Facebook Messenger webhook: GET for verification, POST for incoming messages."""
        VERIFY_TOKEN = env('FACEBOOK_VERIFY_TOKEN', default='magia_fb_webhook_2024')

        if request.method == 'GET':
            mode = request.query_params.get('hub.mode')
            token = request.query_params.get('hub.verify_token')
            challenge = request.query_params.get('hub.challenge')
            if mode == 'subscribe' and token == VERIFY_TOKEN:
                from django.http import HttpResponse
                return HttpResponse(challenge, content_type='text/plain')
            return Response({'error': 'Forbidden'}, status=403)

        # POST: incoming message event
        body = request.data
        if body.get('object') != 'page':
            return Response({'status': 'ignored'})

        for entry in body.get('entry', []):
            page_id = entry.get('id', '')
            config = FacebookConfig.objects.filter(page_id=page_id, is_connected=True).first()
            if not config:
                continue
            user = config.user
            for messaging in entry.get('messaging', []):
                sender_id = messaging.get('sender', {}).get('id', '')
                msg = messaging.get('message', {})
                text = msg.get('text', '').strip()
                if not text or sender_id == page_id:
                    continue  # skip echoes
                # Dedup
                msg_id = msg.get('mid', '')
                if msg_id and ChatMessage.objects.filter(user=user, whatsapp_message_id=msg_id).exists():
                    continue
                # Find active agent with facebook channel
                active_agents = Agent.objects.filter(user=user, is_active=True)
                agent = None
                for a in active_agents:
                    if a.channels and any(str(c).lower() == 'facebook' for c in a.channels):
                        agent = a
                        break
                if not agent:
                    agent = active_agents.first()

                ChatMessage.objects.create(
                    user=user,
                    agent=agent,
                    sender='user',
                    content=text,
                    contact_info=sender_id,
                    source='facebook',
                    whatsapp_message_id=msg_id,
                    status='new'
                )
                # Inbox only - do not auto-add every Messenger sender to CRM
                mark_prospect_replied(user, sender_id, 'facebook')

                # Auto-réponse IA si un agent actif couvre le canal Facebook
                if agent and agent.is_active and agent.channels and any(
                    str(c).lower() == 'facebook' for c in agent.channels
                ):
                    threading.Thread(
                        target=facebook_auto_reply,
                        args=(user.id, agent.id, sender_id, text),
                        daemon=True,
                    ).start()

        return Response({'status': 'EVENT_RECEIVED'})

    @action(detail=True, methods=['post'])
    def sync_messages(self, request, pk=None):
        config = self.get_object()
        if not config.is_connected or not config.page_access_token:
            return Response({'status': 'Facebook non connecté'})
        result = sync_facebook_history(config, request.user)
        if result.get('error'):
            return Response({'error': result['error']}, status=result.get('status_code', 500))
        return Response({'status': 'synced', 'count': result['imported']})

    def perform_destroy(self, instance):
        ChatMessage.objects.filter(user=instance.user, source='facebook').delete()
        Contact.objects.filter(user=instance.user, source='facebook').delete()
        super().perform_destroy(instance)

class AgentTeamViewSet(viewsets.ModelViewSet):
    serializer_class = AgentTeamSerializer
    def get_queryset(self):
        return AgentTeam.objects.filter(user=self.request.user)
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)
        AuditLog.objects.create(
            user=self.request.user,
            action="Création d'équipe",
            details=f"Équipe '{serializer.validated_data['name']}' créée."
        )

    @action(detail=False, methods=['post'])
    def deploy_plan(self, request):
        plan = request.data.get('plan')
        if not plan:
            return Response({'error': 'Plan manquant.'}, status=400)
            
        team = AgentTeam.objects.create(
            user=request.user,
            name=plan['name'],
            description=plan['description'],
            color=plan.get('color', '#1e3a8a')
        )
        
        agents_created = []
        for agent_def in plan['agents']:
            channels = agent_def.get('channels') or infer_channels_for_team_agent(
                agent_def.get('name', ''),
                agent_def.get('role', ''),
                agent_def.get('system_prompt', ''),
            )
            agent = Agent.objects.create(
                user=request.user,
                team=team,
                name=agent_def['name'],
                role=agent_def['role'],
                system_prompt=agent_def['system_prompt'],
                is_team_agent=True,
                is_deployed=True,
                llm_model='gemini-2.0-flash',
                execution_mode='auto',
                channels=channels,
            )
            attach_user_channel_configs(agent, request.user)
            agents_created.append(agent)
            
        for link_def in plan.get('links', []):
            try:
                AgentLink.objects.create(
                    user=request.user,
                    source_agent=agents_created[link_def['from']],
                    target_agent=agents_created[link_def['to']],
                    trigger_type=link_def['trigger'],
                    description=link_def.get('description', '')
                )
            except Exception:
                pass
                
        AuditLog.objects.create(
            user=request.user,
            action="Déploiement d'équipe",
            details=f"Équipe '{team.name}' déployée avec {len(agents_created)} agents."
        )
        
        return Response({'id': team.id, 'status': 'deployed'})

    @action(detail=False, methods=['post'], url_path='design_team')
    def design_team(self, request):
        """
        AI-powered conversational team builder.
        Receives { history: [{role, content}], message: str }
        Returns { reply: str, ready: bool, plan: {...} | null }
        """

        history = request.data.get('history', [])
        user_message = request.data.get('message', '').strip()

        if not user_message:
            return Response({'error': 'Message requis.'}, status=400)

        SYSTEM = """Tu es MAGIA AI, l'architecte expert des solutions d'ingénierie d'agents.
Ton rôle est de concevoir l'équipe de l'utilisateur le plus RAPIDEMENT possible via des choix multiples (QCM).

INSTRUCTIONS DE TON ET STYLE :
- Ton Corporate, Expert et Ultra-Précis.
- NE JAMAIS utiliser de gras (pas de **).
- Réponses épurées.
- Chaque réponse doit être un JSON valide.

STRUCTURE DE RÉPONSE OBLIGATOIRE (JSON) :
{
  "reply": "Ton message court (sans gras)",
  "options": ["Choix 1", "Choix 2", "Choix 3"], 
  "allow_multiple": false,
  "ready": false,
  "plan": null
}

Si plusieurs choix sont possibles, mets "allow_multiple": true.

Si l'équipe est prête :
{
  "reply": "L'équipe est prête.",
  "options": [],
  "ready": true,
  "plan": {
    "name": "Nom prestigieux",
    "description": "Objectif",
    "color": "#1e3a8a",
    "agents": [
      {"name": "Nom Agent", "role": "Expertise", "system_prompt": "Instructions pro sans gras", "channels": ["chat", "whatsapp", "email"]}
    ],
    "links": [
      {"from": 0, "to": 1, "trigger": "interest", "description": "Trigger"}
    ]
  }
}

Triggers : interest, email_requested, whatsapp_requested, manual.
Pour chaque agent, inclus "channels" parmi chat, whatsapp, email selon son rôle."""

        conversation_text = f"SYSTEM: {SYSTEM}\n\n"
        for msg in history[-10:]:
            role = "Utilisateur" if msg.get('role') == 'user' else "MAGIA AI"
            conversation_text += f"{role}: {msg.get('content', '')}\n"
        conversation_text += f"Utilisateur: {user_message}\nMAGIA AI:"

        try:
            api_key = env('GEMINI_API_KEY', default=None)
            if not api_key:
                return Response({'error': 'Clé Gemini manquante dans le fichier .env.'}, status=500)

            client = genai.Client(api_key=api_key)
            
            models_to_try = DEFAULT_GEMINI_MODELS
            reply_text = "Désolé, je rencontre une difficulté."
            
            for m in models_to_try:
                try:
                    response = client.models.generate_content(model=m, contents=conversation_text)
                    reply_text = response.text.strip()
                    break
                except Exception as ex:
                    if m == models_to_try[-1]:
                        return Response({'error': f'Erreur AI : {str(ex)}'}, status=500)
                    continue

            try:
                json_match = re.search(r'(\{.*\})', reply_text, re.DOTALL)
                if json_match:
                    data = json.loads(json_match.group(1))
                    return Response({
                        'reply': data.get('reply', 'Plan conçu.'),
                        'options': data.get('options', []),
                        'ready': data.get('ready', False),
                        'plan': data.get('plan'),
                    })
            except Exception:  
                pass

            return Response({'reply': reply_text, 'options': [], 'ready': False, 'plan': None})

        except Exception as exc:
            logger.error('design_team error: %s\n%s', exc, traceback.format_exc())
            return Response({'error': f'Erreur AI : {exc}'}, status=500)

    @action(detail=True, methods=['post'], url_path='upload_knowledge', parser_classes=[parsers.MultiPartParser, parsers.FormParser])
    def upload_knowledge(self, request, pk=None):
        team = self.get_object()
        file_obj = request.FILES.get('file')
        if not file_obj:
            return Response({'error': 'No file provided'}, status=400)
            
        name = request.data.get('name', getattr(file_obj, 'name', 'Document'))
        
        apply_to_all = request.data.get('apply_to_all_agents', 'false').lower() == 'true'
        agent_ids_str = request.data.get('agent_ids', '')
        
        selected_agent_ids = []
        if apply_to_all:
            selected_agent_ids = list(team.members.values_list('id', flat=True))
        elif agent_ids_str:
            selected_agent_ids = [aid.strip() for aid in agent_ids_str.split(',') if aid.strip()]
            
        try:
            file_binary = file_obj.read()
            file_extension = os.path.splitext(file_obj.name)[1].lower()
            file_obj.seek(0)
            extracted_text = _extract_text(file_obj)
            
            kb_team = KnowledgeBase.objects.create(
                name=name,
                team=team,
                source_type=request.data.get('source_type', 'file'),
                file_binary=file_binary,
                file_extension=file_extension
            )
            
            if extracted_text:
                add_texts_to_knowledge_base(team_id=team.id, raw_text=extracted_text, source_name=kb_team.name)
                
            kbs_created = [{'type': 'team', 'id': kb_team.id}]
            
            for agent_id in selected_agent_ids:
                try:
                    agent = team.members.get(id=agent_id)
                    kb_agent = KnowledgeBase.objects.create(
                        name=name,
                        agent=agent,
                        source_type=request.data.get('source_type', 'file'),
                        file_binary=file_binary,
                        file_extension=file_extension
                    )
                    if extracted_text:
                        add_texts_to_knowledge_base(agent_id=agent.id, raw_text=extracted_text, source_name=kb_agent.name)
                    kbs_created.append({'type': 'agent', 'agent_id': agent.id, 'id': kb_agent.id})
                except Exception:
                    pass
                    
            return Response({'status': 'uploaded', 'results': kbs_created})
        except Exception as exc:
            return Response({'error': str(exc)}, status=400)

class AgentLinkViewSet(viewsets.ModelViewSet):
    serializer_class = AgentLinkSerializer
    def get_queryset(self):
        return AgentLink.objects.filter(user=self.request.user)
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)

class LinkedInConfigViewSet(viewsets.ModelViewSet):
    serializer_class = LinkedInConfigSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        return LinkedInConfig.objects.filter(user=self.request.user)

    def perform_create(self, serializer):
        serializer.save(user=self.request.user)

    @action(detail=True, methods=['get'])
    def refresh_connection(self, request, pk=None):
        return Response({
            "status": "unavailable",
            "message": "LinkedIn n'est plus disponible. Utilisez WhatsApp ou Email.",
        })

    @action(detail=True, methods=['post'])
    def sync_messages(self, request, pk=None):
        return Response({
            "status": "unavailable",
            "message": "LinkedIn n'est plus disponible. Utilisez WhatsApp ou Email.",
        })

    @action(detail=True, methods=['get'])
    def get_connection_url(self, request, pk=None):
        return Response({
            "error": "LinkedIn n'est plus disponible. Utilisez WhatsApp ou Email."
        }, status=status.HTTP_400_BAD_REQUEST)

    @action(detail=True, methods=['post'])
    def prospect(self, request, pk=None):
        return Response({
            "error": "LinkedIn n'est plus disponible. Utilisez WhatsApp ou Email."
        }, status=status.HTTP_400_BAD_REQUEST)

    def create(self, request, *args, **kwargs):
        return Response({
            "error": "LinkedIn n'est plus disponible. Utilisez WhatsApp ou Email."
        }, status=status.HTTP_400_BAD_REQUEST)

class ContactAssignmentViewSet(viewsets.ModelViewSet):
    serializer_class = ContactAssignmentSerializer
    def get_queryset(self):
        return ContactAssignment.objects.filter(user=self.request.user)

class ContactViewSet(viewsets.ModelViewSet):
    from .serializers import ContactSerializer
    serializer_class = ContactSerializer

    def get_queryset(self):
        # One-shot cleanup of auto-imported CRM junk (sync carnet / chat Manual)
        if self.action == 'list':
            purge_passive_crm_contacts(self.request.user)
        return Contact.objects.filter(user=self.request.user).exclude(source='chat')

    def perform_create(self, serializer):
        source = serializer.validated_data.get('source', '')
        if source == 'linkedin':
            from rest_framework.exceptions import ValidationError
            raise ValidationError({
                'source': (
                    "Le canal linkedin n'est plus disponible. "
                    "Utilisez whatsapp, email ou facebook."
                )
            })
        serializer.save(user=self.request.user)

    @action(detail=False, methods=['post'])
    def run_followups(self, request):
        """Manual trigger for follow-up batch (also run by the background scheduler)."""
        from django.core.management import call_command
        from io import StringIO
        out = StringIO()
        call_command('run_followups', stdout=out)
        return Response({'status': 'ok', 'output': out.getvalue()})

    @action(detail=False, methods=['post'], parser_classes=[parsers.MultiPartParser, parsers.FormParser])
    def import_contacts(self, request):
        """
        Import contacts from a CSV or Excel file.
        Expected columns: name, contact_info, source (whatsapp/email/facebook), notes (optional), status (optional)
        """
        file_obj = request.FILES.get('file')
        if not file_obj:
            return Response({'error': 'Aucun fichier fourni.'}, status=400)

        filename = file_obj.name.lower()
        try:
            if filename.endswith('.csv'):
                df = pd.read_csv(file_obj, dtype=str)
            elif filename.endswith(('.xlsx', '.xls')):
                df = pd.read_excel(file_obj, dtype=str)
            else:
                return Response({'error': 'Format non supporté. Utilisez .csv, .xlsx ou .xls'}, status=400)
        except Exception as exc:
            return Response({'error': f'Erreur de lecture du fichier : {str(exc)}'}, status=400)

        # Normalize column names
        df.columns = [c.strip().lower().replace(' ', '_') for c in df.columns]

        # Check required columns
        if 'contact_info' not in df.columns:
            return Response({
                'error': 'Colonne "contact_info" manquante.',
                'columns_found': list(df.columns),
                'hint': 'Colonnes attendues : name, contact_info, source, notes'
            }, status=400)

        created_count = 0
        skipped_count = 0
        errors = []
        valid_statuses = {'new', 'contacted', 'interested', 'ready', 'no'}
        # LinkedIn outreach is retired; Facebook is supported (PSID ou URL de profil)
        valid_sources = {'whatsapp', 'email', 'facebook', 'linkedin', 'chat'}

        for idx, row in df.iterrows():
            contact_info = str(row.get('contact_info', '')).strip()
            if not contact_info or contact_info == 'nan':
                skipped_count += 1
                continue

            name = str(row.get('name', contact_info)).strip()
            if name == 'nan':
                name = contact_info

            source = str(row.get('source', 'whatsapp')).strip().lower()
            if source == 'linkedin':
                errors.append(
                    f"Ligne {idx + 2}: canal linkedin indisponible - contact importé en 'whatsapp'."
                )
                source = 'whatsapp'
            elif source not in valid_sources:
                source = 'whatsapp'

            notes = str(row.get('notes', '')).strip()
            if notes == 'nan':
                notes = ''

            row_status = str(row.get('status', 'new')).strip().lower()
            if row_status not in valid_statuses:
                row_status = 'new'

            try:
                contact, created = Contact.objects.get_or_create(
                    user=request.user,
                    source=source,
                    contact_info=contact_info,
                    defaults={
                        'name': name,
                        'notes': notes,
                        'status': row_status,
                    }
                )
                if created:
                    created_count += 1
                else:
                    skipped_count += 1
            except Exception as exc:
                errors.append({'row': idx + 2, 'contact_info': contact_info, 'error': str(exc)})

        return Response({
            'status': 'done',
            'created': created_count,
            'skipped': skipped_count,
            'errors': errors,
            'total_rows': len(df)
        })

    @action(detail=True, methods=['post'], url_path='contact_via_agent')
    def contact_via_agent(self, request, pk=None):
        """
        Confie ce prospect à un Agent IA qui génère et envoie un premier message
        via WhatsApp, Email ou Facebook (Messenger) selon le canal du contact.
        """
        contact = self.get_object()
        agent_id = request.data.get('agent_id')

        if not agent_id:
            return Response({'error': 'agent_id requis.'}, status=400)

        try:
            agent = Agent.objects.get(id=agent_id, user=request.user)
        except Agent.DoesNotExist:
            return Response({'error': 'Agent introuvable.'}, status=404)

        from .prospect_search_service import send_agent_intro
        sent, message = send_agent_intro(request.user, contact, agent)
        if not sent:
            return Response({'error': message}, status=400)

        return Response({
            'status': 'sent',
            'message': message,
            'agent': agent.name,
            'source': contact.source,
        })


class ChatMessageViewSet(viewsets.ModelViewSet):
    serializer_class = ChatMessageSerializer

    def get_queryset(self):
        return ChatMessage.objects.filter(user=self.request.user).order_by('-created_at')

class TemplateViewSet(viewsets.ReadOnlyModelViewSet):
    serializer_class = TemplateSerializer

    def get_queryset(self):
        return Template.objects.all()

class KnowledgeBaseViewSet(viewsets.ModelViewSet):
    serializer_class = KnowledgeBaseSerializer
    parser_classes = [parsers.MultiPartParser, parsers.FormParser, parsers.JSONParser]

    def get_queryset(self):
        user_agents = Agent.objects.filter(user=self.request.user).values_list('id', flat=True)
        user_teams = AgentTeam.objects.filter(user=self.request.user).values_list('id', flat=True)
        return KnowledgeBase.objects.filter(
            Q(agent_id__in=user_agents) | Q(team_id__in=user_teams)
        )

    def perform_create(self, serializer):
        file_obj = self.request.FILES.get('file')
        instance = serializer.save()

        if file_obj:
            try:
                instance.file_binary = file_obj.read()
                instance.file_extension = os.path.splitext(file_obj.name)[1].lower()
                file_obj.seek(0)

                extracted_text = _extract_text(file_obj)

                if extracted_text or instance.file_binary:
                    instance.save()
                    if extracted_text:
                        if instance.agent:
                            add_texts_to_knowledge_base(agent_id=instance.agent.id, raw_text=extracted_text, source_name=instance.name)
                        elif instance.team:
                            add_texts_to_knowledge_base(team_id=instance.team.id, raw_text=extracted_text, source_name=instance.name)
            except Exception:  # noqa: BLE001
                pass

@method_decorator(csrf_exempt, name='dispatch')
class AgentViewSet(viewsets.ModelViewSet):
    serializer_class = AgentSerializer


    def get_queryset(self):
        user = self.request.user
        owner_agents = Agent.objects.filter(user=user)
        workspaces_as_member = WorkspaceMember.objects.filter(member_user=user).values_list('workspace_owner', flat=True)
        member_agents = Agent.objects.filter(user__in=workspaces_as_member)
        
        return (owner_agents | member_agents).distinct()

    def check_permissions(self, request):
        super().check_permissions(request)
        if hasattr(self, 'detail') and self.detail:
            try:
                obj = self.get_object()
                self.check_workspace_role(request, obj)
            except Exception:
                pass

    def check_workspace_role(self, request, obj):
        if obj.user == request.user:
            return  
        try:
            member = WorkspaceMember.objects.get(workspace_owner=obj.user, member_user=request.user)
            if request.method not in permissions.SAFE_METHODS:
                if member.role != 'editor':
                    raise PermissionDenied("Accès refusé : vous n'avez pas les droits d'édition dans ce workspace.")
        except WorkspaceMember.DoesNotExist:
            raise PermissionDenied("Accès refusé : vous n'êtes pas membre de ce workspace.")

    def perform_create(self, serializer):
        user = self.request.user
        try:
            subscription = user.subscription
            plan = subscription.plan_name
            num_agents_choice = subscription.num_agents
        except Exception:
            plan = 'gratuit'
            num_agents_choice = 2
        
        limits = PLAN_LIMITS.get(plan, PLAN_LIMITS['gratuit'])
        
        if plan == 'pro':
            max_agents = num_agents_choice
        else:
            max_agents = limits['max_agents']
        
        if max_agents is not None:
            current_count = Agent.objects.filter(user=user).count()
            if current_count >= max_agents:
                raise PermissionDenied(
                    f"Limite atteinte : votre abonnement {plan} personnalisé permet {max_agents} agent(s) maximum. "
                    f"Veuillez passer au niveau supérieur pour créer davantage d'agents."
                )
        agent = serializer.save(user=self.request.user)
        if agent.is_team_agent:
            if not agent.channels:
                agent.channels = infer_channels_for_team_agent(
                    agent.name, agent.role, agent.system_prompt
                )
                agent.save(update_fields=['channels'])
            attach_user_channel_configs(agent, self.request.user)

    @action(detail=False, methods=['get'])
    def dashboard_stats(self, request):
        user = request.user
        agents = Agent.objects.filter(user=user)
        agent_ids = agents.values_list('id', flat=True)

        messages = ChatMessage.objects.filter(agent_id__in=agent_ids)
        human_conversations = messages.filter(sender='user').count()
        ai_responses = messages.filter(sender='ai').count()
        
        pertinent_escalations = messages.filter(status='pertinent').count()
        automation_rate = "100%"
        if human_conversations > 0:
            rate = max(0, 100 - (pertinent_escalations / human_conversations * 100))
            automation_rate = f"{int(rate)}%"

        total_minutes_saved = ai_responses * 5
        hours_saved = total_minutes_saved // 60
        time_saved_str = f"{hours_saved}h" if hours_saved > 0 else f"{total_minutes_saved}m"

        economy = hours_saved * 15
        active_agents_count = agents.filter(is_deployed=True).count()

        recent_msgs = messages.filter(sender='ai').select_related('agent').order_by('-created_at')[:5]
        
        now = timezone.now()
        recent_activities = []
        colors = ['bg-blue-500', 'bg-blue-400', 'bg-blue-600', 'bg-blue-900', 'bg-blue-300']
        
        for i, msg in enumerate(recent_msgs):
            diff = now - msg.created_at
            if diff.total_seconds() < 60:
                time_str = "À l'instant"
            elif diff.total_seconds() < 3600:
                time_str = f"{int(diff.total_seconds() // 60)}m"
            elif diff.total_seconds() < 86400:
                time_str = f"{int(diff.total_seconds() // 3600)}h"
            else:
                time_str = f"{int(diff.total_seconds() // 86400)}j"

            recent_activities.append({
                'user': msg.agent.name,
                'action': msg.content[:30] + '...' if len(msg.content) > 30 else msg.content,
                'color': colors[i % len(colors)],
                'time': time_str,
                'avatar': msg.agent.avatar if msg.agent.avatar else None
            })

        # Aggregation pour le graphique (derniers 7 jours)
        last_7_days_cutoff = timezone.now() - datetime.timedelta(days=7)
        daily_stats = (
            messages.filter(created_at__gte=last_7_days_cutoff)
            .annotate(day=TruncDay('created_at'))
            .values('day', 'sender')
            .annotate(count=Count('id'))
            .order_by('day')
        )
        
        graph_dict = {}
        for i in range(7):
            d = (timezone.now() - datetime.timedelta(days=6-i)).date()
            graph_dict[d.isoformat()] = {
                "name": d.strftime('%d/%m'), 
                "conversations": 0, 
                "responses": 0
            }
            
        for entry in daily_stats:
            d_str = entry['day'].date().isoformat()
            if d_str in graph_dict:
                if entry['sender'] == 'user':
                    graph_dict[d_str]['conversations'] = entry['count']
                else:
                    graph_dict[d_str]['responses'] = entry['count']
        
        graph_data = list(graph_dict.values())

        return Response({
            'active_agents': active_agents_count,
            'conversations': human_conversations,
            'ai_responses': ai_responses,
            'automation_rate': automation_rate,
            'time_saved': time_saved_str,
            'economy_eur': f"{economy} €",
            'csat': "94.8%",
            'recent_activity': recent_activities,
            'graph_data': graph_data
        })

    @action(detail=True, methods=['post'])
    def feedback(self, request, pk=None):
        agent = self.get_object()
        serializer = AgentFeedbackSerializer(data=request.data)
        if serializer.is_valid():
            serializer.save(user=request.user, agent=agent)
            return Response(serializer.data, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
        


    @action(detail=True, methods=['post'], url_path='upload_knowledge', parser_classes=[parsers.MultiPartParser, parsers.FormParser])
    def upload_knowledge(self, request, pk=None):
        agent = self.get_object()
        file_obj = request.FILES.get('file')
        if not file_obj:
            return Response({'error': 'No file provided'}, status=400)
            
        name = request.data.get('name', getattr(file_obj, 'name', 'Document'))
        kb = KnowledgeBase(
            name=name,
            agent=agent,
            source_type=request.data.get('source_type', 'file'),
        )
        try:
            kb.file_binary = file_obj.read()
            kb.file_extension = os.path.splitext(file_obj.name)[1].lower()
            file_obj.seek(0)

            extracted_text = _extract_text(file_obj)

            if extracted_text or kb.file_binary:
                kb.save()
                if extracted_text:
                    add_texts_to_knowledge_base(agent.id, extracted_text, kb.name)
            return Response({'status': 'uploaded', 'kb_id': kb.id})
        except Exception as exc:
            return Response({'error': str(exc)}, status=400)

    @action(detail=True, methods=['post'])
    def deploy(self, request, pk=None):
        agent = self.get_object()
        if not agent.knowledge_bases.exists():
            return Response(
                {'error': "Cet agent ne peut pas être déployé sans base de connaissances. Veuillez en ajouter au moins une."},
                status=status.HTTP_400_BAD_REQUEST
            )
        time.sleep(1) 
        agent.is_deployed = True
        agent.save()
        return Response({'status': 'deployed'})

    @action(detail=True, methods=['post'])
    def toggle_pause(self, request, pk=None):
        agent = self.get_object()
        new_status = not agent.is_active
        agent.is_active = new_status
        agent.save()

        if not new_status:
            self._cascade_pause(agent)

        return Response({'is_active': agent.is_active})

    def _cascade_pause(self, agent):
        """Recursively pause all agents that are targets of links from this agent."""
        links = AgentLink.objects.filter(source_agent=agent)
        for link in links:
            target = link.target_agent
            if target.is_active:
                target.is_active = False
                target.save()
                self._cascade_pause(target)

    @action(detail=True, methods=['post'])
    def sandbox_chat(self, request, pk=None):
        agent = self.get_object()
        user_message = request.data.get('message', '')
        
        team_id = agent.team.id if agent.team else None
        rag_context = search_agent_and_team_knowledge_base(agent_id=agent.id, team_id=team_id, query=user_message, top_k=4)
        context_for_llm = f"Documents sources (Extraits pertinents par recherche sémantique RAG) :\n{rag_context if rag_context.strip() else '[Aucun document fourni]'}\n\n[MODE SANDBOX : Ce message est un test isolé]"
        
        user_plan = 'gratuit'
        if hasattr(request.user, 'subscription'):
            user_plan = request.user.subscription.plan_name

        allowed, reason = check_ai_quota(request.user)
        if not allowed:
            return Response({'error': reason}, status=403)

        response = get_llm_response(
            agent_name=agent.name,
            agent_role=agent.role,
            system_prompt=agent.system_prompt,
            knowledge_context=context_for_llm,
            user_message=user_message,
            model_name=agent.llm_model,
            user_plan=user_plan
        )

        return Response({
            'reply': response,
            'agent': agent.name,
            'timestamp': time.time()
        })

    def _get_assigned_agent(self, user, contact_info, default_agent):
        if not contact_info:
            return default_agent
        assignment = ContactAssignment.objects.filter(user=user, contact_info=contact_info).first()
        return assignment.agent if assignment else default_agent

    def _check_handoff(self, agent, user_message, contact_info, source=None):
        """
        Check if any trigger is met to hand off to another agent.
        On success: reassign contact, log + optionally send an intro from the new agent.
        Returns the target agent or None.
        """
        links = AgentLink.objects.filter(source_agent=agent)
        if not links.exists():
            return None

        from .llm_service import classify_handoff
        for link in links:
            # "manual" is reserved for human / UI escalation - skip auto LLM classification
            if link.trigger_type == 'manual':
                continue
            if classify_handoff(link.trigger_type, user_message):
                target = link.target_agent
                ContactAssignment.objects.update_or_create(
                    user=agent.user,
                    contact_info=contact_info,
                    defaults={'agent': target},
                )

                intro = build_handoff_intro(target, previous_agent=agent)
                resolved_source = source
                if not resolved_source:
                    contact = Contact.objects.filter(
                        user=agent.user, contact_info=contact_info
                    ).first()
                    resolved_source = contact.source if contact else 'chat'

                ChatMessage.objects.create(
                    user=agent.user,
                    agent=target,
                    sender='ai',
                    content=intro,
                    contact_info=contact_info,
                    source=resolved_source or 'chat',
                    status='new',
                )

                if resolved_source in ('whatsapp', 'email'):
                    try:
                        MessagingService.send_message(
                            agent.user, contact_info, intro, resolved_source
                        )
                    except Exception as exc:
                        logger.warning("Handoff intro send failed: %s", exc)

                AuditLog.objects.create(
                    user=agent.user,
                    action="Handoff d'agent",
                    details=(
                        f"{agent.name} → {target.name} "
                        f"(trigger={link.trigger_type}, contact={contact_info})"
                    ),
                )
                return target
        return None

    @action(detail=True, methods=['post'])
    def chat(self, request, pk=None):
        agent = self.get_object()
        if not agent.is_active:
            return Response({'status': 'ignored', 'reason': 'agent_paused', 'reply': "Cette unité IA est actuellement en pause et ne peut pas répondre."})
        
        user_message = request.data.get('message', '')
        contact_info = request.data.get('contact_info', 'anonymous')
        
        effective_agent = self._get_assigned_agent(request.user, contact_info, agent)
        
        handed_to = self._check_handoff(effective_agent, user_message, contact_info, source='chat')
        if handed_to:
            effective_agent = handed_to

        mark_prospect_replied(request.user, contact_info)
        
        ChatMessage.objects.create(agent=effective_agent, user=request.user, sender='user', content=user_message, contact_info=contact_info)

        history = list(ChatMessage.objects.filter(agent=effective_agent, contact_info=contact_info).order_by('-created_at')[:6])
        history.reverse()
        history_str = "\n".join([f"{m.sender}: {m.content}" for m in history])

        rag_context = search_knowledge_base(agent.id, user_message, top_k=4)
        context_for_llm = f"Documents sources (Extraits pertinents par recherche sémantique RAG) :\n{rag_context if rag_context.strip() else '[Aucun document fourni]'}\n\nHistorique récent :\n{history_str}"
        user_plan = 'gratuit'
        if hasattr(request.user, 'subscription'):
            user_plan = request.user.subscription.plan_name

        allowed, reason = check_ai_quota(request.user)
        if not allowed:
            return Response({'status': 'ignored', 'reason': 'quota_reached', 'reply': reason})

        response = get_llm_response(
            agent_name=effective_agent.name,
            agent_role=effective_agent.role,
            system_prompt=effective_agent.system_prompt,
            knowledge_context=context_for_llm,
            user_message=user_message,
            model_name=effective_agent.llm_model,
            user_plan=user_plan
        )

        ChatMessage.objects.create(agent=effective_agent, user=request.user, sender='ai', content=response, contact_info=contact_info)
        schedule_followup_after_ai(request.user, contact_info, analyze=True)
        return Response({
            'reply': response,
            'agent': effective_agent.name,
            'handoff': handed_to.name if handed_to else None,
            'timestamp': time.time()
        })

    @action(detail=True, methods=['post'])
    def email_process(self, request, pk=None):
        process_emails_for_agent(pk)
        return Response({'status': 'processed'})

    @action(detail=False, methods=['post'])
    def process_all_emails(self, request):
        import logging
        logger = logging.getLogger('agents.email_service')
        agents = Agent.objects.all()
        count = 0
        logger.info(f"Déclenchement du traitement global des emails pour {agents.count()} agents potential.")
        for agent in agents:
            has_email = False
            if agent.channels:
                has_email = any(str(c).lower() == 'email' for c in agent.channels)
            
            if agent.is_active and has_email:
                logger.info(f"Traitement des emails pour l'agent: {agent.name} (ID: {agent.id})")
                process_emails_for_agent(agent.id)
                count += 1
            else:
                if not agent.is_active:
                    logger.debug(f"Agent {agent.name} ignoré car inactif.")
                if not has_email:
                    logger.debug(f"Agent {agent.name} ignoré car canal email non activé.")
        
        return Response({'status': 'processed', 'agent_count': count})

    @action(detail=True, methods=['post'])
    def whatsapp_process(self, request, pk=None):
        agent = self.get_object()
        if not agent.is_active:
            return Response({'status': 'ignored', 'reason': 'agent_paused'})
            
        message = request.data.get('message', '').strip()
        wa_id = request.data.get('message_id', '')
        wa_sender = request.data.get('sender', '')

        effective_agent = self._get_assigned_agent(request.user, wa_sender, agent)
        
        handed_to = self._check_handoff(effective_agent, message, wa_sender, source='whatsapp')
        if handed_to:
            effective_agent = handed_to

        mark_prospect_replied(request.user, wa_sender, 'whatsapp')

        config = effective_agent.whatsapp_config
        is_connected = config is not None and config.is_connected
        if not is_connected and not request.data.get('force'):
            return Response({'status': 'ignored', 'reason': 'no_whatsapp_account_connected'})

        if not message:
            return Response({'status': 'ignored', 'reason': 'empty_message'})

        status = classify_pertinence(effective_agent.role, message)

        rag_context = search_knowledge_base(effective_agent.id, message, top_k=4)
        context_for_llm = f"Documents sources (Extraits pertinents par recherche sémantique RAG) :\n{rag_context if rag_context.strip() else '[Aucun document fourni]'}"
        
        answer = get_llm_response(
            agent_name=effective_agent.name,
            agent_role=effective_agent.role,
            system_prompt=effective_agent.system_prompt,
            knowledge_context=context_for_llm,
            user_message=message,
            model_name=effective_agent.llm_model
        )

        ChatMessage.objects.create(agent=effective_agent, user=request.user, sender='user', content=message, contact_info=wa_sender, source='whatsapp', whatsapp_message_id=wa_id, status=status)
        ChatMessage.objects.create(agent=effective_agent, user=request.user, sender='ai', content=answer, contact_info=wa_sender, source='whatsapp', status='new')
        schedule_followup_after_ai(request.user, wa_sender, 'whatsapp', analyze=True)

        return Response({
            'status': 'responded',
            'response': answer,
            'source': 'ai',
            'agent': effective_agent.name,
            'handoff': handed_to.name if handed_to else None,
        })

    @action(detail=True, methods=['get'])
    def prospections(self, request, pk=None):
        agent = self.get_object()
        messages = ChatMessage.objects.filter(agent=agent).exclude(contact_info__isnull=True).order_by('created_at')
        threads = {}
        for m in messages:
            if m.contact_info not in threads:
                threads[m.contact_info] = {
                    'contact': m.contact_info,
                    'source': m.source,
                    'messages': [],
                    'last_updated': m.created_at
                }
            threads[m.contact_info]['messages'].append({
                'id': m.id,
                'sender': m.sender,
                'content': m.content,
                'created_at': m.created_at
            })
            threads[m.contact_info]['last_updated'] = m.created_at
        
        result = list(threads.values())
        result.sort(key=lambda x: x['last_updated'], reverse=True)
        return Response(result)

    @action(detail=True, methods=['post'])
    def send_manual_reply(self, request, pk=None):
        agent = self.get_object()
        contact_info = request.data.get('contact_info')
        content = request.data.get('content')
        source = request.data.get('source')
        
        email_config_id = request.data.get('email_config_id')
        
        if not contact_info or not content or not source:
            return Response({'error': 'Missing data'}, status=400)
            
        if source == 'chat':
            # Chatbot / widget thread only - never mixes WhatsApp/Email history
            ChatMessage.objects.create(
                agent=agent,
                user=request.user,
                sender='user',
                content=content,
                contact_info=contact_info,
                source=source,
                status='new',
            )
            history = list(
                ChatMessage.objects.filter(
                    agent=agent, source='chat', contact_info=contact_info
                ).order_by('-created_at')[:6]
            )
            history.reverse()
            history_str = "\n".join([f"{m.sender}: {m.content}" for m in history])

            rag_context = search_knowledge_base(agent.id, content, top_k=4)
            context_for_llm = f"Documents sources (Extraits pertinents par recherche sémantique RAG) :\n{rag_context if rag_context.strip() else '[Aucun document fourni]'}\n\nHistorique récent :\n{history_str}"

            response_text = get_llm_response(
                agent_name=agent.name,
                agent_role=agent.role,
                system_prompt=agent.system_prompt,
                knowledge_context=context_for_llm,
                user_message=content,
                model_name=agent.llm_model
            )
            ChatMessage.objects.create(
                agent=agent,
                user=request.user,
                sender='ai',
                content=response_text,
                contact_info=contact_info,
                source=source,
                status='new',
            )
            schedule_followup_after_ai(request.user, contact_info, source, analyze=False)
            return Response({'status': 'sent', 'reply': response_text})

        elif source == 'email':
            config = None
            if email_config_id:
                config = EmailConfig.objects.filter(id=email_config_id, user=request.user).first()

            if not config and agent and agent.email_config:
                config = agent.email_config

            if not config:
                config = EmailConfig.objects.filter(user=request.user).filter(
                    Q(is_active=True) | Q(oauth_token__isnull=False)
                ).first()

            if not config:
                return Response({'error': 'Reliez votre email avec MAGIA pour effectuer cette opération'}, status=400)

            success = send_email_reply(config, contact_info, "Re: Contact", content)
            if not success:
                return Response({'error': 'Échec de l\'envoi email. Vérifiez la connexion dans Paramètres.'}, status=500)
            ChatMessage.objects.create(
                agent=agent,
                user=request.user,
                sender='ai',
                content=content,
                contact_info=contact_info,
                source=source,
                status='new',
                email_config=config,
            )
            schedule_followup_after_ai(request.user, contact_info, source, analyze=False)

        elif source == 'whatsapp':
            success, wa_err = MessagingService.send_whatsapp_detailed(
                request.user, contact_info, content
            )
            if not success:
                return Response({
                    'error': wa_err or (
                        'WhatsApp non connecté ou service indisponible. '
                        'Activez WhatsApp dans Paramètres.'
                    )
                }, status=500)
            ChatMessage.objects.create(
                agent=agent,
                user=request.user,
                sender='ai',
                content=content,
                contact_info=contact_info,
                source=source,
                status='new',
            )
            schedule_followup_after_ai(request.user, contact_info, source, analyze=False)

        elif source == 'facebook':
            success, fb_err = MessagingService.send_facebook_detailed(
                request.user, contact_info, content
            )
            if not success:
                return Response({
                    'error': fb_err or (
                        "Échec de l'envoi Messenger. Vérifiez que la Page est connectée et que "
                        "le contact a écrit à la Page il y a moins de 24h (règle Meta)."
                    )
                }, status=500)
            ChatMessage.objects.create(
                agent=agent,
                user=request.user,
                sender='ai',
                content=content,
                contact_info=contact_info,
                source=source,
                status='new',
            )
            schedule_followup_after_ai(request.user, contact_info, source, analyze=False)

        elif source == 'linkedin':
            return Response({
                'error': "Le canal linkedin n'est plus disponible. Utilisez WhatsApp, Email ou Facebook."
            }, status=400)
                
        return Response({'status': 'sent'})
        
    @action(detail=False, methods=['post'])
    def universal_reply(self, request):
        contact_info = request.data.get('contact_info')
        content = request.data.get('content')
        source = request.data.get('source')
        email_config_id = request.data.get('email_config_id')
        
        if not contact_info or not content or not source:
            return Response({'error': 'Missing data'}, status=400)
            
        if source == 'email':
            config = None
            if email_config_id:
                config = EmailConfig.objects.filter(id=email_config_id, user=request.user).first()
            if not config:
                config = EmailConfig.objects.filter(user=request.user).first()
                
            if not config:
                return Response({'error': 'Reliez votre email avec MAGIA pour effectuer cette opération'}, status=400)
                
            success = send_email_reply(config, contact_info, "Re: Contact", content)
            if not success:
                return Response({'error': 'Failed to send email'}, status=500)
            
            ChatMessage.objects.create(
                user=request.user,
                sender='ai',
                content=content,
                contact_info=contact_info,
                source=source,
                status='pertinent',
                email_config=config
            )
            schedule_followup_after_ai(request.user, contact_info, source, analyze=False)
            return Response({'status': 'sent'})
            
        elif source == 'whatsapp':
            success, wa_err = MessagingService.send_whatsapp_detailed(
                request.user, contact_info, content
            )
            if success:
                ChatMessage.objects.create(
                    user=request.user,
                    sender='ai',
                    content=content,
                    contact_info=contact_info,
                    source=source,
                    status='pertinent',
                )
                schedule_followup_after_ai(request.user, contact_info, source, analyze=False)
                return Response({'status': 'sent'})
            return Response({
                'error': wa_err or 'Échec de l\'envoi WhatsApp. Reconnectez WhatsApp dans Paramètres.'
            }, status=500)

        elif source == 'facebook':
            success, fb_err = MessagingService.send_facebook_detailed(
                request.user, contact_info, content
            )
            if success:
                ChatMessage.objects.create(
                    user=request.user,
                    sender='ai',
                    content=content,
                    contact_info=contact_info,
                    source=source,
                    status='pertinent',
                )
                schedule_followup_after_ai(request.user, contact_info, source, analyze=False)
                return Response({'status': 'sent'})
            return Response({
                'error': fb_err or (
                    "Échec de l'envoi Messenger. Vérifiez que la Page est connectée et que "
                    "le contact a écrit à la Page il y a moins de 24h (règle Meta)."
                )
            }, status=500)

        elif source == 'linkedin':
            return Response({
                'error': "Le canal linkedin n'est plus disponible. Utilisez WhatsApp, Email ou Facebook."
            }, status=400)

        return Response({'error': 'Source not supported for direct reply'}, status=400)

    @action(detail=True, methods=['get'])
    def messages(self, request, pk=None):
        """
        Agent chatbot threads only by default (source=chat).
        Channel traffic (whatsapp/email/facebook) belongs in inbox channel tabs.
        """
        agent = self.get_object()
        source = request.query_params.get('source', 'chat')
        contact = request.query_params.get('contact')
        msgs = ChatMessage.objects.filter(agent=agent)
        if source:
            msgs = msgs.filter(source=source)
        if contact:
            msgs = msgs.filter(contact_info=contact)
        msgs = msgs.order_by('created_at')
        serializer = ChatMessageSerializer(msgs, many=True)
        return Response(serializer.data)

    @action(detail=False, methods=['post'])
    def sync_inbox(self, request):
        """
        Re-sync email / Facebook (and optionally restart WhatsApp session) so
        the inbox shows a fuller message history.
        Body: { "source": "email" | "whatsapp" | "facebook" | "all" }
        """
        source = (request.data.get('source') or 'all').lower()
        results = {'email': None, 'whatsapp': None, 'facebook': None}

        if source in ('email', 'all'):
            configs = EmailConfig.objects.filter(user=request.user).filter(
                Q(is_active=True) | Q(oauth_token__isnull=False)
            )
            if not configs.exists():
                results['email'] = {
                    'status': 'error',
                    'error': 'Aucun compte email connecté. Reliez Gmail dans Paramètres.',
                }
            else:
                total = {'imported': 0, 'skipped': 0, 'errors': []}
                for config in configs:
                    # Ensure active for sync path
                    if not config.is_active and config.oauth_token:
                        config.is_active = True
                        config.save(update_fields=['is_active'])
                    stats = sync_email_history(config)
                    total['imported'] += stats.get('imported', 0)
                    total['skipped'] += stats.get('skipped', 0)
                    if stats.get('error'):
                        total['errors'].append(stats['error'])
                results['email'] = {'status': 'done', **total}

        if source in ('whatsapp', 'all'):
            from .whatsapp_process import start_for_user, is_running
            wa = WhatsAppConfig.objects.filter(user=request.user).first()
            if not wa:
                results['whatsapp'] = {
                    'status': 'error',
                    'error': 'Aucun WhatsApp configuré. Connectez WhatsApp dans Paramètres.',
                }
            else:
                auth = request.META.get('HTTP_AUTHORIZATION', '')
                token = auth.split(' ', 1)[1].strip() if auth.lower().startswith('bearer ') else ''
                start = start_for_user(request.user, auth_token=token)
                results['whatsapp'] = {
                    'status': 'running' if is_running(request.user.id) or start.get('already_running') or start.get('started') else 'error',
                    'message': (
                        'Service WhatsApp actif - l\'historique continue de se synchroniser '
                        'en arrière-plan (texte + médias indiqués).'
                    ),
                    'error': start.get('error'),
                    'started': start.get('started'),
                    'already_running': start.get('already_running'),
                }

        if source in ('facebook', 'all'):
            fb = FacebookConfig.objects.filter(user=request.user, is_connected=True).first()
            if not fb or not fb.page_access_token:
                results['facebook'] = {
                    'status': 'error',
                    'error': 'Aucune Page Facebook connectée. Connectez-la dans Paramètres.',
                }
            else:
                stats = sync_facebook_history(fb, request.user)
                if stats.get('error'):
                    results['facebook'] = {'status': 'error', 'error': stats['error']}
                else:
                    results['facebook'] = {'status': 'done', 'imported': stats['imported']}

        return Response(results)

    @action(detail=False, methods=['get'])
    def all_conversations(self, request):
        user_agents = Agent.objects.filter(user=request.user).values_list('id', flat=True)
        search = request.query_params.get('search', '')
        source_filter = request.query_params.get('source', '')

        # Build base queryset
        base_qs = ChatMessage.objects.filter(
            Q(user=request.user) | Q(agent_id__in=user_agents)
        )
        if source_filter:
            base_qs = base_qs.filter(source=source_filter)
        else:
            # "Tout" = channel inboxes only; chatbot threads live under Agents
            base_qs = base_qs.exclude(source='chat')

        # Search across contacts
        if search:
            # Also search Contact table names
            matched_contacts = Contact.objects.filter(
                user=request.user,
                name__icontains=search
            ).values_list('contact_info', flat=True)
            base_qs = base_qs.filter(
                Q(contact_info__icontains=search) |
                Q(contact_name__icontains=search) |
                Q(contact_info__in=matched_contacts)
            )

        # Get latest message ID for each unique (source, contact_info)
        latest_msg_ids = base_qs.values('source', 'contact_info').annotate(
            latest_id=MaxAgg('id')
        ).values_list('latest_id', flat=True)

        messages = ChatMessage.objects.filter(
            id__in=latest_msg_ids
        ).select_related('agent').order_by('-created_at')

        # Build a lookup for Contact names
        contact_info_list = [m.contact_info for m in messages if m.contact_info]
        contacts_lookup = {
            (c.source, c.contact_info): c.name
            for c in Contact.objects.filter(
                user=request.user,
                contact_info__in=contact_info_list
            )
        }

        result = []
        for m in messages:
            agent_data = None
            if m.agent:
                agent_data = {'id': str(m.agent.id), 'name': m.agent.name}

            # Best name: Contact table > message.contact_name > raw contact_info
            display_name = (
                contacts_lookup.get((m.source, m.contact_info))
                or m.contact_name
                or m.contact_info
            )

            unread = ChatMessage.objects.filter(
                Q(user=request.user) | Q(agent_id__in=user_agents),
                source=m.source,
                contact_info=m.contact_info,
                is_read=False,
                sender='user'
            ).count()

            result.append({
                'contact': m.contact_info,
                'contact_name': display_name,
                'avatar_letter': (display_name or '?')[0].upper(),
                'source': m.source,
                'agent': agent_data,
                'email_config_id': (
                    m.email_config_id
                    if m.email_config_id
                    else (m.agent.email_config_id if m.agent and m.agent.email_config_id else None)
                ),
                'unread': unread,
                'lastMsg': {
                    'id': m.id,
                    'content': m.content[:120],
                    'sender': m.sender,
                    'created_at': m.created_at,
                    'is_me': m.sender == 'ai',
                },
                'last_updated': m.created_at
            })

        return Response(result)

    @action(detail=False, methods=['get'])
    def contact_messages(self, request):
        contact = request.query_params.get('contact')
        source = request.query_params.get('source')
        if not contact or not source:
            return Response({'error': 'Missing contact or source parameters'}, status=400)

        user_agents = list(Agent.objects.filter(user=request.user).values_list('id', flat=True))
        # Return ALL messages (no limit) so user can scroll back years
        messages = ChatMessage.objects.filter(
            Q(user=request.user) | Q(agent_id__in=user_agents),
            contact_info=contact,
            source=source
        ).order_by('created_at').select_related('agent')

        # Mark messages as read
        messages.filter(is_read=False, sender='user').update(is_read=True)

        serializer = ChatMessageSerializer(messages, many=True)
        return Response(serializer.data)

    @action(detail=False, methods=['get'])
    def contacts_list(self, request):
        """
        Inbox search/autocomplete from message history (not CRM sync dump).
        Falls back to intentional CRM contacts for numbers not yet messaged in MAGIA.
        """
        search = request.query_params.get('search', '')
        source = request.query_params.get('source', '')
        user_agents = Agent.objects.filter(user=request.user).values_list('id', flat=True)

        msg_qs = ChatMessage.objects.filter(
            Q(user=request.user) | Q(agent_id__in=user_agents)
        ).exclude(contact_info__isnull=True).exclude(contact_info='')
        if source:
            msg_qs = msg_qs.filter(source=source)
        if search:
            msg_qs = msg_qs.filter(
                Q(contact_info__icontains=search) | Q(contact_name__icontains=search)
            )

        latest_ids = msg_qs.values('source', 'contact_info').annotate(
            latest_id=MaxAgg('id')
        ).values_list('latest_id', flat=True)
        latest_msgs = ChatMessage.objects.filter(id__in=latest_ids).order_by('-created_at')[:100]

        seen = set()
        result = []
        for m in latest_msgs:
            key = (m.source, m.contact_info)
            if key in seen:
                continue
            seen.add(key)
            name = m.contact_name or m.contact_info
            result.append({
                'id': m.id,
                'name': name,
                'contact_info': m.contact_info,
                'source': m.source,
                'avatar_letter': (name or '?')[0].upper(),
                'last_message_at': m.created_at,
            })

        crm = Contact.objects.filter(user=request.user).exclude(source='chat')
        if source:
            crm = crm.filter(source=source)
        if search:
            crm = crm.filter(
                Q(name__icontains=search) | Q(contact_info__icontains=search)
            )
        for c in crm.order_by('-last_message_at', '-created_at')[:50]:
            key = (c.source, c.contact_info)
            if key in seen:
                continue
            seen.add(key)
            name = c.name or c.contact_info
            result.append({
                'id': c.id,
                'name': name,
                'contact_info': c.contact_info,
                'source': c.source,
                'avatar_letter': (name or '?')[0].upper(),
                'last_message_at': c.last_message_at,
            })

        return Response(result[:100])

    @action(detail=False, methods=['get'])
    def boite_reception(self, request):
        user_agents = Agent.objects.filter(user=request.user).values_list('id', flat=True)
        messages = ChatMessage.objects.filter(
            agent_id__in=user_agents, 
            status='pertinent'
        ).order_by('-created_at')
        
        serializer = ChatMessageSerializer(messages, many=True)
        return Response(serializer.data)

class AuditLogViewSet(viewsets.ReadOnlyModelViewSet):
    serializer_class = AuditLogSerializer
    def get_queryset(self):
        return AuditLog.objects.filter(user=self.request.user).order_by('-created_at')


class ProspectSearchJobViewSet(viewsets.ModelViewSet):
    """
    Recherche Apollo → enrichissement → CRM → envoi auto Email/WhatsApp
    (+ ajout CRM des profils Facebook pour prise de contact manuelle).
    """
    serializer_class = ProspectSearchJobSerializer
    http_method_names = ['get', 'post', 'head', 'options']

    def get_queryset(self):
        return (
            ProspectSearchJob.objects.filter(user=self.request.user)
            .prefetch_related('leads')
            .select_related('agent')
            .order_by('-created_at')
        )

    def create(self, request, *args, **kwargs):
        from django.conf import settings as dj_settings
        from .prospect_search_service import MAX_RESULTS_CAP, start_job_async

        if not (getattr(dj_settings, 'APOLLO_API_KEY', None) or '').strip():
            return Response(
                {
                    'error': (
                        "APOLLO_API_KEY manquante. Ajoutez votre clé Apollo dans backend/.env "
                        "puis redémarrez Django."
                    )
                },
                status=400,
            )

        agent_id = request.data.get('agent_id')
        if not agent_id:
            return Response({'error': 'agent_id requis.'}, status=400)

        try:
            agent = Agent.objects.get(id=agent_id, user=request.user)
        except Agent.DoesNotExist:
            return Response({'error': 'Agent introuvable.'}, status=404)

        if not agent.is_deployed:
            return Response({'error': "L'agent doit être déployé."}, status=400)

        channels = (request.data.get('channels') or 'both').lower()
        if channels not in ('email', 'whatsapp', 'facebook', 'both', 'all'):
            return Response(
                {'error': 'channels doit être email, whatsapp, facebook, both ou all.'},
                status=400,
            )

        try:
            max_results = int(request.data.get('max_results') or 10)
        except (TypeError, ValueError):
            max_results = 10
        max_results = min(max(max_results, 1), MAX_RESULTS_CAP)

        filters = request.data.get('filters') or {}
        if not isinstance(filters, dict):
            return Response({'error': 'filters doit être un objet JSON.'}, status=400)

        # Normalize simple string filters from the UI
        for key in (
            'person_titles', 'person_locations', 'organization_locations',
            'q_organization_keyword_tags', 'person_seniorities',
        ):
            if key in filters and isinstance(filters[key], str):
                filters[key] = [t.strip() for t in filters[key].split(',') if t.strip()]

        job = ProspectSearchJob.objects.create(
            user=request.user,
            agent=agent,
            filters=filters,
            channels=channels,
            max_results=max_results,
            status='pending',
        )

        from django.db import transaction as db_transaction

        def _kick():
            start_job_async(job.id)

        db_transaction.on_commit(_kick)

        serializer = self.get_serializer(job)
        return Response(serializer.data, status=201)


@method_decorator(csrf_exempt, name='dispatch')
class ApolloPhoneWebhookView(viewsets.ViewSet):
    """Apollo async phone enrichment callback."""
    permission_classes = [AllowAny]
    authentication_classes = []

    def create(self, request):
        from django.conf import settings as dj_settings
        from .prospect_search_service import handle_apollo_phone_webhook

        expected = (getattr(dj_settings, 'APOLLO_WEBHOOK_SECRET', None) or '').strip()
        token = (
            request.query_params.get('token')
            or request.headers.get('X-Apollo-Token')
            or ''
        ).strip()
        if expected and token != expected:
            return Response({'error': 'Unauthorized'}, status=401)

        job_id = request.query_params.get('job_id')
        lead_id = request.query_params.get('lead_id')
        try:
            job_id_int = int(job_id) if job_id else None
        except ValueError:
            job_id_int = None
        try:
            lead_id_int = int(lead_id) if lead_id else None
        except ValueError:
            lead_id_int = None

        payload = request.data if isinstance(request.data, dict) else {}
        result = handle_apollo_phone_webhook(payload, job_id=job_id_int, lead_id=lead_id_int)
        return Response({'status': 'ok', **result})

